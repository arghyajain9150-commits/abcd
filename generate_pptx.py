import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("python-pptx not yet ready. Please ensure pip install python-pptx has finished.")
    sys.exit(1)

def create_champ_presentation(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette: Deep Forest Green, Emerald, Amber Gold, Pure White, Soft Charcoal
    COLOR_BG = RGBColor(23, 50, 44)        # Deep Forest #17322C
    COLOR_PRIMARY = RGBColor(47, 122, 104) # Emerald Green #2F7A68
    COLOR_ACCENT = RGBColor(227, 165, 66)  # Amber Gold #E3A542
    COLOR_TEXT_LIGHT = RGBColor(245, 247, 243)
    COLOR_TEXT_MUTED = RGBColor(163, 217, 201)
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_DARK = RGBColor(23, 50, 44)
    COLOR_CARD_BG = RGBColor(245, 247, 243)
    COLOR_CARD_BORDER = RGBColor(225, 227, 218)

    blank_layout = prs.slide_layouts[6]

    def add_header(slide, title_text, category_text="CHAMP · CAMPUS HEALTH & APPOINTMENT MANAGEMENT PLATFORM"):
        # Top Category Banner
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT

        # Main Slide Title
        tx_box2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf2 = tx_box2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_DARK

    # ─────────────────────────────────────────────────────────────
    # SLIDE 1: Title Slide (Dark Luxury Theme)
    # ─────────────────────────────────────────────────────────────
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_BG
    bg1.line.fill.background()

    # Title Text Frame
    tb = slide1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "CHAMP"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    p_sub = tf.add_paragraph()
    p_sub.text = "Campus Health & Appointment Management Platform"
    p_sub.font.size = Pt(24)
    p_sub.font.bold = True
    p_sub.font.color.rgb = COLOR_ACCENT
    p_sub.space_before = Pt(10)

    p_desc = tf.add_paragraph()
    p_desc.text = "An Autonomous, AI-Powered Healthcare Ecosystem with Micro-Spatial Outbreak Radar, Vision AI Rx OCR, 2FA Pharmacy Dispensary & Open Innovation Data Hub"
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = COLOR_TEXT_MUTED
    p_desc.space_before = Pt(15)

    p_track = tf.add_paragraph()
    p_track.text = "🚀 Tracks: AI / Machine Learning & Open Innovation | IIT / University Healthcare Hackathon"
    p_track.font.size = Pt(12)
    p_track.font.bold = True
    p_track.font.color.rgb = COLOR_WHITE
    p_track.space_before = Pt(35)

    # ─────────────────────────────────────────────────────────────
    # SLIDE 2: Problem Statement & Campus Reality
    # ─────────────────────────────────────────────────────────────
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "The Critical Problem in Campus Healthcare Today", "01. The Problem Space")

    cards_data = [
        ("High-Density Contagion Velocity", "Universities and hostels house thousands of students in close proximity. A single case of eye flu (conjunctivitis) or viral fever explodes across entire hostel floors before the medical center even notices.", "🦠 Outbreak Lag"),
        ("Fragmented OPD Queue Bottlenecks", "Students wait 2+ hours in crowded waiting rooms while sick, exacerbating transmission. Doctors lack digital patient queuing, leading to chaotic walk-in traffic and administrative burnout.", "⏳ Queue Paralysis"),
        ("Prescription Leakage & Drug Insecurity", "Paper prescriptions get lost, leading to unauthorized medicine pickup, zero allergy cross-checks, and no visibility into campus dispensary inventory levels or stockouts.", "💊 Paper Inefficiency")
    ]

    for i, (title, desc, badge) in enumerate(cards_data):
        left = Inches(0.8 + i * 3.95)
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(3.8), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CARD_BORDER

        tb = slide2.shapes.add_textbox(left + Inches(0.2), Inches(2.0), Inches(3.4), Inches(4.4))
        tf = tb.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = badge
        p0.font.size = Pt(11)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_PRIMARY

        p1 = tf.add_paragraph()
        p1.text = title
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_DARK
        p1.space_before = Pt(8)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(91, 113, 105)
        p2.space_before = Pt(12)

    # ─────────────────────────────────────────────────────────────
    # SLIDE 3: The CHAMP Solution & Architectural Innovation
    # ─────────────────────────────────────────────────────────────
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "CHAMP: The 360° Connected Campus Healthcare System", "02. The Solution")

    sol_cards = [
        ("🎓 1. Student Health Portal", "Instant doctor appointment booking with conflict detection, 2FA digital medicine pickup pass, Gemini AI health triage, lab report vault, and wellness community forum.", COLOR_PRIMARY),
        ("🩺 2. Multi-Desk OPD Portal", "Dedicated doctor consultation queues, clinical allergy cross-checks, live formulary typeahead autocomplete, and direct AI Outbreak Radar diagnostic tagging.", COLOR_PRIMARY),
        ("💊 3. Smart Dispensary Desk", "Real-time prescription fulfillment, 2FA 4-digit OTP / QR handover verification, atomic inventory stock decrementing, and automated restock alerts.", COLOR_PRIMARY),
        ("🌐 4. Open Data & AI Engine", "7-Day Outbreak ML Projection with R₀ What-If transmission simulator, Gemini Vision handwritten Rx OCR scanner, and public REST Developer APIs.", COLOR_ACCENT)
    ]

    for i, (title, desc, color) in enumerate(sol_cards):
        row = i // 2
        col = i % 2
        left = Inches(0.8 + col * 5.95)
        top = Inches(1.8 + row * 2.5)

        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.8), Inches(2.3))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CARD_BORDER

        tb = slide3.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), Inches(5.3), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = color

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(91, 113, 105)
        p2.space_before = Pt(8)

    # ─────────────────────────────────────────────────────────────
    # SLIDE 4: AI/ML Track Innovations
    # ─────────────────────────────────────────────────────────────
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "AI/ML Track: Predictive Outbreak Intelligence & Vision OCR", "03. AI/ML Innovations")

    ai_features = [
        ("📈 7-Day Outbreak ML Projection & R₀ Simulator", "• Ingests live doctor contagious tags (e.g. Conjunctivitis, Flu) and maps micro-spatial hostel clusters.\n• Computes effective Reproduction Rate (R₀ = 1.84) and generates dynamic 7-day infection curves.\n• Interactive What-If Simulator: Judges can test intervention sliders (0–100% compliance) to watch projected peaks drop in real time."),
        ("📷 Gemini Vision AI Prescription Digitizer", "• Solves the notorious 'unreadable doctor handwriting' problem on campus.\n• Students or dispensary staff upload handwritten slips; Vision AI segments cursive notes into structured medication objects with confidence scoring.\n• 1-Click action immediately creates a validated digital pharmacy order."),
        ("🧠 Smart Gemini 2.5 Symptom Triage", "• Interactive conversational clinical triage with prompt chips ('Itchy eye', 'Fever').\n• Evaluates clinical urgency (Low/Moderate/Severe) and deep-links directly into the exact recommended medical specialist desk (e.g. Dr. Sanjana Iyer, Dermatology).")
    ]

    for i, (title, text) in enumerate(ai_features):
        left = Inches(0.8 + i * 3.95)
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(3.8), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CARD_BORDER

        tb = slide4.shapes.add_textbox(left + Inches(0.2), Inches(2.0), Inches(3.4), Inches(4.4))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_PRIMARY

        p2 = tf.add_paragraph()
        p2.text = text
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = RGBColor(91, 113, 105)
        p2.space_before = Pt(10)

    # ─────────────────────────────────────────────────────────────
    # SLIDE 5: Open Innovation Track Features
    # ─────────────────────────────────────────────────────────────
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "Open Innovation Track: Public Health Interoperability", "04. Open Innovation")

    open_features = [
        ("🌐 1. Public Campus Health REST APIs (/open-data)", "• Provides open, standards-compliant REST endpoints for student developers and civic hackers:\n   - GET /api/open/stats (Live anonymized aggregate metrics & bed occupancy)\n   - GET /api/open/outbreaks (Geospatial cluster feeds for university epidemiologists)\n   - GET /api/open/pharmacy-stock (Essential drug availability index)\n• Complete with interactive API Playground, JSON viewer, and cURL / Python / JS snippets under CC-BY 4.0 Open License."),
        ("🛡️ 2. Privacy-Preserving Research Consent Engine", "• Students can toggle 'Opt-in to Anonymized Research Data Sharing' directly in their portal.\n• Enables federated campus disease research while guaranteeing ZERO Personally Identifiable Information (PII) exposure (NDHM & HIPAA principles)."),
        ("🤝 3. Crowd-Sourced Student Wellness Forum", "• Community-driven peer support hub where students publish and upvote actionable health initiatives (Hostel ORS Hydration drives, 20-20-20 screen rule, Exam isolation meal buddies).")
    ]

    for i, (title, text) in enumerate(open_features):
        top = Inches(1.8 + i * 1.65)
        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.7), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CARD_BORDER

        tb = slide5.shapes.add_textbox(Inches(1.0), top + Inches(0.15), Inches(11.3), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(13.5)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_PRIMARY

        p2 = tf.add_paragraph()
        p2.text = text
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(91, 113, 105)
        p2.space_before = Pt(4)

    # ─────────────────────────────────────────────────────────────
    # SLIDE 6: Production-Grade Security & Clinical Safety Architecture
    # ─────────────────────────────────────────────────────────────
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "Production-Grade Security & Clinical Safeguards", "05. Security & Safety")

    sec_cards = [
        ("🔐 2FA Medicine Handover (OTP & QR)", "Prescriptions generate a dynamic 4-digit OTP & scannable QR with 7-day expiration. Pharmacists must verify the OTP before dispensing, preventing wrongful medication release.", "Security"),
        ("⚛️ Atomic Inventory Transactions", "Stock updates execute within strict database transaction locks, preventing race conditions or phantom stock during high-volume dispensary rushes.", "Data Integrity"),
        ("⚠️ Clinical Allergy Cross-Check", "Doctor Prescription Writer scans patient allergy records (Penicillin, Sulfa) in real time and displays high-visibility warning banners prior to Rx issuance.", "Clinical Safety"),
        ("🖨️ Clean A4 @media Print Isolation", "Dedicated print stylesheet isolates only the official stamped prescription slip, stripping all backdrop overlays, mockups, and UI artifacts for physical medical records.", "Compliance")
    ]

    for i, (title, desc, cat) in enumerate(sec_cards):
        row = i // 2
        col = i % 2
        left = Inches(0.8 + col * 5.95)
        top = Inches(1.8 + row * 2.5)

        card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.8), Inches(2.3))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CARD_BORDER

        tb = slide6.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), Inches(5.3), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = cat.upper()
        p0.font.size = Pt(10)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_ACCENT

        p1 = tf.add_paragraph()
        p1.text = title
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_DARK
        p1.space_before = Pt(4)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = RGBColor(91, 113, 105)
        p2.space_before = Pt(6)

    # ─────────────────────────────────────────────────────────────
    # SLIDE 7: Technical Stack & Live Deployment
    # ─────────────────────────────────────────────────────────────
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "Technical Architecture & Live Production Stack", "06. Architecture")

    tech_data = [
        ("Frontend Layer", "React 19 · Vite 8 · TanStack Query · React Router v7 · Lucide Icons · Responsive 1160px Desktop Dashboard + PWA Mobile Viewport", COLOR_PRIMARY),
        ("Backend Layer", "Node.js (ES Modules) · Express.js · Socket.io (Real-time Queue & Radar events) · In-Memory Rate Limiting · RateLimiter Middleware", COLOR_PRIMARY),
        ("AI / ML Engines", "Google Gemini 2.5 API (Clinical Triage) · Gemini Vision OCR (Handwritten Rx Parser) · Custom Spatial-Temporal Clustering & R₀ Transmission Engine", COLOR_PRIMARY),
        ("Database & Cloud", "PostgreSQL via Supabase Cloud Pooler · Vercel Production Frontend · Render Cloud Backend · GitHub Continuous Deployment", COLOR_ACCENT)
    ]

    for i, (title, desc, col) in enumerate(tech_data):
        top = Inches(1.8 + i * 1.25)
        card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.7), Inches(1.1))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CARD_BORDER

        tb = slide7.shapes.add_textbox(Inches(1.0), top + Inches(0.12), Inches(11.3), Inches(0.85))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = col

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(91, 113, 105)
        p2.space_before = Pt(4)

    # ─────────────────────────────────────────────────────────────
    # SLIDE 8: Summary & Conclusion
    # ─────────────────────────────────────────────────────────────
    slide8 = prs.slides.add_slide(blank_layout)
    bg8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg8.fill.solid()
    bg8.fill.fore_color.rgb = COLOR_BG
    bg8.line.fill.background()

    tb = slide8.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Transforming Campus Healthcare into a Predictive, Safe Ecosystem"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    p1 = tf.add_paragraph()
    p1.text = "CHAMP is not just a booking app — it is an active epidemiological containment network, an AI diagnostic assistant, and an open innovation research platform for universities worldwide."
    p1.font.size = Pt(16)
    p1.font.color.rgb = COLOR_TEXT_MUTED
    p1.space_before = Pt(16)

    p2 = tf.add_paragraph()
    p2.text = "🌐 Live URL: https://abcd-five-zeta.vercel.app  |  💻 GitHub: github.com/arghyajain9150-commits/abcd"
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_ACCENT
    p2.space_before = Pt(30)

    prs.save(output_path)
    print("PowerPoint generated successfully at: " + output_path)

if __name__ == "__main__":
    desktop_dir = os.path.expanduser('~/OneDrive/Desktop')
    if not os.path.exists(desktop_dir):
        desktop_dir = os.path.expanduser('~/Desktop')
    out_file = os.path.join(desktop_dir, 'CHAMP_Hackathon_Pitch_Presentation.pptx')
    create_champ_presentation(out_file)
